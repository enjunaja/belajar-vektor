pip install python-pptx

import streamlit as st
from pptx import Presentation
from io import BytesIO

def coba_menu():
    # Judul aplikasi
    st.markdown("<h3 style='text-align: center;'> Ini menu untuk coba program </h3>", unsafe_allow_html=True)
    st.markdown("<h6 style='text-align: center;'> 💻 Media Pembelajaran Ruang Vektor Berbasis DNR!</h3>", unsafe_allow_html=True)
    st.write("")
    st.markdown("<h4 style='text-align: left;'>Silahkan tulis kode programnya</h4>", unsafe_allow_html=True)
    


    # Mendownload file PowerPoint dari URL (ganti URL dengan yang sesuai)
    url_pptx = 'https://github.com/enjunaja/belajar-vektor/blob/main/Organisasi%20Objek%20Matematika%20Ruang%20Vektor.pptx'
    response = requests.get(url_pptx)
    pptx_data = BytesIO(response.content)

    # Memuat presentasi dari file PowerPoint
    presentation = Presentation(pptx_data)

    # Menampilkan slide presentasi di Streamlit
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                st.write(shape.text)


    
